import autogen
import requests
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# OpenRouter API configuration
OPENROUTER_API_KEY = os.environ['OPENROUTER_API_KEY']

class OpenRouterLLM:
    def __init__(self, model="openai/gpt-4o-2024-08-06"):
        self.model = model

    def create_completion(self, messages):
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}"
            },
            data=json.dumps({
                "model": self.model,
                "messages": messages,
                "max_tokens": 4000,
                "temperature": 0.7,
            })
        )
        return response.json()

llm = OpenRouterLLM()

content_strategist_config = {
    "name": "ContentStrategist",
    "system_message": """You are a presentation content strategist.
    Your responsibilities:
    - Create detailed presentation outlines
    - Define key messages and takeaways
    - Structure content flow and progression
    - Ensure content alignment with presentation goals
    - Recommend content distribution across slides"""
}

slide_designer_config = {
    "name": "SlideDesigner",
    "system_message": """You are a slide design specialist.
    Your responsibilities:
    - Create visually appealing slide layouts
    - Recommend color schemes and typography
    - Design information hierarchy on slides
    - Suggest visual elements and imagery
    - Ensure consistent design language"""
}

content_writer_config = {
    "name": "ContentWriter",
    "system_message": """You are a presentation content writer.
    Your responsibilities:
    - Write clear and concise slide content
    - Create engaging headlines and bullet points
    - Develop speaker notes
    - Ensure content clarity and impact
    - Maintain consistent tone and style"""
}

class PresentationProgress:
    def __init__(self):
        self.prs = Presentation()
        self.slides = []
        self.total_slides = 0
        self.target_slides = 8

    def add_slide(self, content, layout=None):
        if layout is None:
            layout = self.prs.slide_layouts[1]  # Use title and content layout by default

        slide = self.prs.slides.add_slide(layout)
        self.slides.append({
            'content': content,
            'slide_object': slide
        })
        self.total_slides += 1

    def apply_content_to_slide(self, slide_dict):
        slide = slide_dict['slide_object']
        content = slide_dict['content']

        # Apply title
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            if 'title' in content:
                slide.shapes.title.text = content['title']

        # Apply content
        if 'body' in content:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape != slide.shapes.title:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        p.text = content['body']

        # Apply speaker notes if they exist
        if 'notes' in content:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content['notes']

    def save(self, filename="presentation.pptx"):
        self.prs.save(filename)

    def is_complete(self):
        return self.total_slides >= self.target_slides

    def get_progress(self):
        return f"Current progress: {self.total_slides}/{self.target_slides} slides"

presentation_progress = PresentationProgress()

def create_messages(system_message, user_message):
    return [
        {"role": "system", "content": system_message},
        {"role": "user", "content": f"{user_message}\n\nCurrent progress: {presentation_progress.get_progress()}"}
    ]

class CustomAssistantAgent(autogen.AssistantAgent):
    def generate_reply(self, sender=None, messages=None):
        if messages is None:
            messages = self._oai_messages

        if not messages:
            return None

        system_message = self.system_message

        last_message = messages[-1]
        if isinstance(last_message, dict):
            last_message_content = last_message.get('content', '')
        else:
            last_message_content = str(last_message)

        api_messages = create_messages(system_message, last_message_content)
        response = llm.create_completion(api_messages)

        try:
            reply = response['choices'][0]['message']['content']
            if self.name == "content_writer":
                try:
                    content_dict = json.loads(reply)
                    presentation_progress.add_slide(content_dict)
                except json.JSONDecodeError:
                    print("Warning: Could not parse content writer response as JSON")
            return reply
        except KeyError:
            return "I apologize, but I encountered an error processing your request."

content_strategist = CustomAssistantAgent(
    name="content_strategist",
    system_message=content_strategist_config["system_message"]
)

slide_designer = CustomAssistantAgent(
    name="slide_designer",
    system_message=slide_designer_config["system_message"]
)

content_writer = CustomAssistantAgent(
    name="content_writer",
    system_message=content_writer_config["system_message"]
)

user_proxy = autogen.UserProxyAgent(
    name="user_proxy",
    human_input_mode="TERMINATE",
    max_consecutive_auto_reply=3
)

groupchat = autogen.GroupChat(
    agents=[user_proxy, content_strategist, slide_designer, content_writer],
    messages=[],
    max_round=15
)

manager = autogen.GroupChatManager(groupchat=groupchat)

initial_message = """Create a professional presentation about artificial intelligence trends in 2024.
Requirements:
1. At least 8 slides
2. Include executive summary
3. Cover key trends, market analysis, and future predictions
4. Provide speaker notes for each slide

Please start by creating a detailed outline and design guidelines."""

# Continue generating content until slide count is met
while not presentation_progress.is_complete():
    if presentation_progress.total_slides == 0:
        user_proxy.initiate_chat(manager, message=initial_message)
    else:
        next_slide_prompt = f"""Create the next slide content and design.
Current progress: {presentation_progress.get_progress()}
Ensure content is clear, impactful, and visually appealing."""
        user_proxy.initiate_chat(manager, message=next_slide_prompt)

# Apply all content to slides and save
for slide in presentation_progress.slides:
    presentation_progress.apply_content_to_slide(slide)
presentation_progress.save()

print(f"\nFinal presentation statistics:\n{presentation_progress.get_progress()}")
