import autogen
import requests
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import time

# OpenRouter API configuration
OPENROUTER_API_KEY = os.environ['OPENROUTER_API_KEY']

class OpenRouterLLM:
    def __init__(self, model="openai/gpt-4o-2024-08-06"):
        self.model = model
        
    def create_completion(self, messages):
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}"
            },
            data=json.dumps({
                "model": self.model,
                "messages": messages,
                "max_tokens": 4000,
                "temperature": 0.7,
            })
        )
        return response.json()

class PresentationBuilder:
    def __init__(self, filename="ai_trends_2024.pptx"):
        self.prs = Presentation()
        self.filename = filename
        self.slide_count = 0
        self.target_slides = 15
        
        # Set default slide size to widescreen
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Define theme colors
        self.colors = {
            'primary': RGBColor(33, 60, 114),
            'secondary': RGBColor(0, 133, 202),
            'accent': RGBColor(242, 144, 0),
            'text': RGBColor(51, 51, 51),
            'background': RGBColor(255, 255, 255)
        }
        
    def create_title_slide(self, title, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self.slide_count += 1
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.color.rgb = self.colors['primary']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = self.colors['secondary']
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def create_content_slide(self, content):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.slide_count += 1
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = content['title']
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = self.colors['primary']
        title_para.font.bold = True
        
        # Add body content
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(content['body']):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['text']
            p.level = 0
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = content['notes']
        
    def create_chart_slide(self, title, chart_data, chart_type="bar"):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = self.colors['primary']
        title_para.font.bold = True
        
        # Add placeholder for chart
        chart_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1.5),
            Inches(8),
            Inches(5)
        )
        chart_placeholder.fill.solid()
        chart_placeholder.fill.fore_color.rgb = self.colors['secondary']
        
        # Add caption
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_frame.text = "Data visualization placeholder"
        caption_para = caption_frame.paragraphs[0]
        caption_para.font.size = Pt(14)
        caption_para.font.color.rgb = self.colors['text']
        caption_para.alignment = PP_ALIGN.CENTER
    
    def save(self):
        timestamp = time.strftime("%Y%m%d-%H%M%S")
        version_filename = f"{self.filename.rsplit('.', 1)[0]}_{timestamp}.pptx"
        self.prs.save(version_filename)
        self.prs.save(self.filename)
        return version_filename
    
    def is_complete(self):
        return self.slide_count >= self.target_slides
    
    def get_progress(self):
        return f"Current progress: {self.slide_count}/{self.target_slides} slides"

# Agent configurations
content_strategist_config = {
    "name": "ContentStrategist",
    "system_message": """You are a presentation content strategist.
    Your responsibilities:
    - Create detailed presentation outlines
    - Define key messages and takeaways
    - Structure content flow and progression
    - Ensure content alignment with presentation goals
    - Recommend content distribution across slides
    
    For the AI trends presentation, focus on:
    1. Market size and growth
    2. Key technology developments
    3. Industry adoption rates
    4. Future predictions
    5. Challenges and opportunities"""
}

content_writer_config = {
    "name": "ContentWriter",
    "system_message": """You are a presentation content writer who MUST ALWAYS return content in valid JSON format.

    Your response format MUST be valid JSON like this:
    {
        "title": "Slide Title Here",
        "body": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
        "notes": "Speaker notes here"
    }

    Your responsibilities:
    - Write clear and concise slide content
    - Create engaging headlines and bullet points
    - Develop speaker notes
    - Ensure content clarity and impact
    - Maintain consistent tone and style
    
    CRITICAL: 
    1. Always wrap your entire response in valid JSON
    2. Use double quotes for strings
    3. Use array format for bullet points
    4. Include all three keys: title, body, notes
    5. Do not include any text outside the JSON structure"""
}

visualization_expert_config = {
    "name": "VisualizationExpert",
    "system_message": """You are a data visualization specialist.
    Your responsibilities:
    - Suggest appropriate chart types
    - Structure data for visual representation
    - Create clear and impactful visualizations
    - Ensure data storytelling
    - Recommend visualization layouts
    
    Return your suggestions in JSON format:
    {
        "title": "Chart Title",
        "type": "bar|line|pie",
        "data": {"label1": value1, "label2": value2}
    }"""
}

class CustomAssistantAgent(autogen.AssistantAgent):
    def __init__(self, name, system_message, presentation_builder):
        super().__init__(name=name, system_message=system_message)
        self.presentation_builder = presentation_builder
        self.llm = OpenRouterLLM()
    
    def generate_reply(self, messages=None, sender=None, config=None):
        """
        Generate a reply based on the messages and handle slide creation.
        
        Args:
            messages: The messages to respond to
            sender: The sender of the message
            config: Additional configuration
            
        Returns:
            str: The generated reply
        """
        if messages is None:
            messages = self._oai_messages
            
        if not messages:
            return None
            
        # Get the last message content
        last_message = messages[-1]
        last_message_content = last_message.content if hasattr(last_message, 'content') else str(last_message)
            
        api_messages = [
            {"role": "system", "content": self.system_message},
            {"role": "user", "content": last_message_content}
        ]
        
        response = self.llm.create_completion(api_messages)
        
        try:
            reply = response['choices'][0]['message']['content']
            
            # Handle content writer responses
            if self.name == "content_writer":
                try:
                    content_dict = json.loads(reply)
                    self.presentation_builder.create_content_slide(content_dict)
                    print(f"\nAdded new slide: {content_dict['title']}")
                except json.JSONDecodeError as e:
                    print(f"Warning: JSON parsing error: {str(e)}")
                    default_content = {
                        "title": "Content Processing Error",
                        "body": ["Content could not be properly formatted", 
                               "Please review and update this slide"],
                        "notes": "This slide needs to be reviewed and updated with proper content."
                    }
                    self.presentation_builder.create_content_slide(default_content)
                    
            # Handle visualization expert responses
            elif self.name == "visualization_expert":
                try:
                    viz_dict = json.loads(reply)
                    self.presentation_builder.create_chart_slide(
                        viz_dict['title'],
                        viz_dict['data'],
                        viz_dict['type']
                    )
                    print(f"\nAdded new chart slide: {viz_dict['title']}")
                except json.JSONDecodeError as e:
                    print(f"Warning: JSON parsing error: {str(e)}")
                    default_chart = {
                        "title": "Visualization Error",
                        "data": {"Error": 100},
                        "type": "bar"
                    }
                    self.presentation_builder.create_chart_slide(**default_chart)
            
            return reply
        except Exception as e:
            print(f"Error generating reply: {str(e)}")
            return "I apologize, but I encountered an error processing your request."

def generate_presentation():
    # Initialize presentation builder
    builder = PresentationBuilder()
    
    # Create title slide
    builder.create_title_slide(
        "Artificial Intelligence Trends 2024",
        "A Comprehensive Market Analysis and Future Outlook"
    )
    
    # Initialize agents
    content_strategist = CustomAssistantAgent(
        name="content_strategist",
        system_message=content_strategist_config["system_message"],
        presentation_builder=builder
    )
    
    content_writer = CustomAssistantAgent(
        name="content_writer",
        system_message=content_writer_config["system_message"],
        presentation_builder=builder
    )
    
    visualization_expert = CustomAssistantAgent(
        name="visualization_expert",
        system_message=visualization_expert_config["system_message"],
        presentation_builder=builder
    )
    
    user_proxy = autogen.UserProxyAgent(
        name="user_proxy",
        human_input_mode="TERMINATE",
        max_consecutive_auto_reply=10
    )
    
    # Create group chat
    groupchat = autogen.GroupChat(
        agents=[user_proxy, content_strategist, content_writer, visualization_expert],
        messages=[],
        max_round=50
    )
    
    manager = autogen.GroupChatManager(groupchat=groupchat)
    
    # Generate presentation content
    initial_message = """Create a professional presentation about artificial intelligence trends in 2024.
    Requirements:
    1. Include executive summary
    2. Cover key trends, market analysis, and future predictions
    3. Include data visualizations and charts
    4. Provide speaker notes for each slide"""
    
    while not builder.is_complete():
        if builder.slide_count <= 1:  # Only title slide exists
            user_proxy.initiate_chat(manager, message=initial_message)
        else:
            next_slide_prompt = f"""Create the next slide content and design.
            {builder.get_progress()}
            Ensure content flows naturally from the previous slides."""
            user_proxy.initiate_chat(manager, message=next_slide_prompt)
    
    # Save final presentation
    final_file = builder.save()
    print(f"\nPresentation completed: {final_file}")
    print(builder.get_progress())
    
    return final_file

if __name__ == "__main__":
    
    # Generate the presentation
    output_file = generate_presentation()
    print(f"Presentation generated successfully: {output_file}")

