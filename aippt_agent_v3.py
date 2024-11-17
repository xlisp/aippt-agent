import autogen
import requests
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import time

# OpenRouter API configuration
OPENROUTER_API_KEY = os.environ['OPENROUTER_API_KEY']

class OpenRouterLLM:
    def __init__(self, model="openai/gpt-4o-2024-08-06"):
        self.model = model
        
    def create_completion(self, messages):
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}"
            },
            data=json.dumps({
                "model": self.model,
                "messages": messages,
                "max_tokens": 4000,
                "temperature": 0.7,
            })
        )
        return response.json()

llm = OpenRouterLLM()

content_strategist_config = {
    "name": "ContentStrategist",
    "system_message": """You are a presentation content strategist.
    Generate content in JSON format with keys: slide_number, title, body, notes.
    Focus on one slide at a time with clear, impactful content."""
}

slide_designer_config = {
    "name": "SlideDesigner",
    "system_message": """You are a slide design specialist.
    Generate design instructions in JSON format with keys: layout_type, colors, fonts, spacing.
    Provide specific layout recommendations for each slide."""
}

content_writer_config = {
    "name": "ContentWriter",
    "system_message": """You are a presentation content writer.
    Generate final slide content in JSON format with keys: title, body, notes.
    Each response should be complete content for one slide."""
}

class PresentationProgress:
    def __init__(self):
        self.filename = "presentation_v3.pptx"
        self.total_slides = 0
        self.target_slides = 5
        self.current_outline = None
        
    def initialize_presentation(self):
        self.prs = Presentation()
        
    def add_slide(self, content):
        """Add a single slide and immediately save the presentation"""
        try:
            # Create new presentation if it doesn't exist
            if not hasattr(self, 'prs'):
                self.initialize_presentation()
            
            # Add slide
            layout = self.prs.slide_layouts[1]  # Using title and content layout
            slide = self.prs.slides.add_slide(layout)
            
            # Apply content
            if 'title' in content:
                slide.shapes.title.text = content['title']
            
            if 'body' in content:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        if shape != slide.shapes.title:
                            text_frame = shape.text_frame
                            text_frame.clear()
                            p = text_frame.paragraphs[0]
                            p.text = content['body']
            
            # Add speaker notes
            if 'notes' in content:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = content['notes']
            
            self.total_slides += 1
            
            # Save after each slide
            self.save()
            
            print(f"Added slide {self.total_slides} and saved presentation")
            return True
            
        except Exception as e:
            print(f"Error adding slide: {str(e)}")
            return False
    
    def save(self):
        """Save the presentation with a timestamp"""
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        save_filename = f"presentation_{timestamp}.pptx"
        self.prs.save(save_filename)
        
        # Also save as the main presentation file
        self.prs.save(self.filename)
    
    def is_complete(self):
        return self.total_slides >= self.target_slides
    
    def get_progress(self):
        return f"Current progress: {self.total_slides}/{self.target_slides} slides"

presentation_progress = PresentationProgress()

def create_messages(system_message, user_message):
    return [
        {"role": "system", "content": system_message},
        {"role": "user", "content": f"{user_message}\n\nCurrent progress: {presentation_progress.get_progress()}"}
    ]

class CustomAssistantAgent(autogen.AssistantAgent):
    def generate_reply(self, sender=None, messages=None):
        if messages is None:
            messages = self._oai_messages
            
        if not messages:
            return None
            
        system_message = self.system_message
        
        last_message = messages[-1]
        if isinstance(last_message, dict):
            last_message_content = last_message.get('content', '')
        else:
            last_message_content = str(last_message)
        
        api_messages = create_messages(system_message, last_message_content)
        response = llm.create_completion(api_messages)
        
        try:
            reply = response['choices'][0]['message']['content']
            if self.name == "content_writer":
                try:
                    content_dict = json.loads(reply)
                    presentation_progress.add_slide(content_dict)
                except json.JSONDecodeError:
                    print("Warning: Could not parse content writer response as JSON")
            return reply
        except KeyError:
            return "I apologize, but I encountered an error processing your request."

content_strategist = CustomAssistantAgent(
    name="content_strategist",
    system_message=content_strategist_config["system_message"]
)

slide_designer = CustomAssistantAgent(
    name="slide_designer",
    system_message=slide_designer_config["system_message"]
)

content_writer = CustomAssistantAgent(
    name="content_writer",
    system_message=content_writer_config["system_message"]
)

user_proxy = autogen.UserProxyAgent(
    name="user_proxy",
    human_input_mode="TERMINATE",
    max_consecutive_auto_reply=3  # Reduced number of consecutive replies
)

groupchat = autogen.GroupChat(
    agents=[user_proxy, content_strategist, slide_designer, content_writer],
    messages=[],
    max_round=15  # Reduced number of rounds
)

manager = autogen.GroupChatManager(groupchat=groupchat)

def generate_presentation(topic):
    """Generate the presentation one slide at a time"""
    presentation_progress.initialize_presentation()
    
    while not presentation_progress.is_complete():
        current_slide = presentation_progress.total_slides + 1
        
        slide_prompt = f"""Create slide {current_slide} for the presentation about {topic}.
        Current progress: {presentation_progress.get_progress()}
        Generate complete content for this single slide, including title, body, and speaker notes.
        Return in JSON format."""
        
        user_proxy.initiate_chat(
            manager,
            message=slide_prompt,
            max_consecutive_auto_reply=3
        )
        
        # Save point check
        if presentation_progress.total_slides % 3 == 0:
            print(f"\nCheckpoint saved at {presentation_progress.get_progress()}")
        
        # Small delay between slides to prevent rate limiting
        time.sleep(2)

# Main execution
if __name__ == "__main__":
    presentation_topic = "Artificial Intelligence Trends in 2024"
    generate_presentation(presentation_topic)
    print(f"\nFinal presentation statistics:\n{presentation_progress.get_progress()}")

