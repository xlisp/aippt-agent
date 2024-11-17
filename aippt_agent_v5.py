import autogen
import requests
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import time

# OpenRouter API configuration
OPENROUTER_API_KEY = os.environ['OPENROUTER_API_KEY']

class OpenRouterLLM:
    def __init__(self, model="openai/gpt-4o-2024-08-06"):
        self.model = model
        
    def create_completion(self, messages):
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}"
            },
            data=json.dumps({
                "model": self.model,
                "messages": messages,
                "max_tokens": 4000,
                "temperature": 0.7,
            })
        )
        return response.json()

llm = OpenRouterLLM()

content_strategist_config = {
    "name": "ContentStrategist",
    "system_message": """You are a presentation content strategist.
    Your responsibilities:
    - Create detailed presentation outlines
    - Define key messages and takeaways
    - Structure content flow and progression
    - Ensure content alignment with presentation goals
    - Recommend content distribution across slides"""
}

slide_designer_config = {
    "name": "SlideDesigner",
    "system_message": """You are a slide design specialist.
    Your responsibilities:
    - Create visually appealing slide layouts
    - Recommend color schemes and typography
    - Design information hierarchy on slides
    - Suggest visual elements and imagery
    - Ensure consistent design language"""
}

content_writer_config = {
    "name": "ContentWriter",
    "system_message": """You are a presentation content writer who MUST ALWAYS return content in valid JSON format.

    Your response format MUST be valid JSON like this:
    {
        "title": "Slide Title Here",
        "body": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
        "notes": "Speaker notes here"
    }

    Your responsibilities:
    - Write clear and concise slide content
    - Create engaging headlines and bullet points
    - Develop speaker notes
    - Ensure content clarity and impact
    - Maintain consistent tone and style
    
    CRITICAL: 
    1. Always wrap your entire response in valid JSON
    2. Use double quotes for strings
    3. Use array format for bullet points
    4. Include all three keys: title, body, notes
    5. Do not include any text outside the JSON structure
    6. Do not include markdown or other formatting

    Example response:
    {
        "title": "AI Market Growth 2024",
        "body": [
            "Global AI market reached $120B in 2024",
            "37% YoY growth in enterprise adoption",
            "Key sectors: Healthcare, Finance, Manufacturing"
        ],
        "notes": "This slide highlights the remarkable growth in AI adoption across industries, with particular emphasis on the three key sectors that showed the highest integration rates."
    }"""
}

visualization_expert_config = {
    "name": "VisualizationExpert",
    "system_message": """You are a data visualization specialist.
    Your responsibilities:
    - Suggest appropriate chart types
    - Structure data for visual representation
    - Create clear and impactful visualizations
    - Ensure data storytelling
    - Recommend visualization layouts"""
}

class PresentationProgress:
    def __init__(self, filename="presentation.pptx"):
        self.filename = filename
        self.prs = Presentation()
        self.slides = []
        self.total_slides = 0
        self.target_slides = 15
        self.version = 0
        
    def add_slide(self, content, layout=None):
        # Validate content structure
        if not isinstance(content, dict):
            print("Warning: Content must be a dictionary")
            return
        
        # Ensure all required keys are present
        required_keys = ['title', 'body', 'notes']
        for key in required_keys:
            if key not in content:
                print(f"Warning: Missing required key '{key}' in slide content")
                if key == 'body':
                    content[key] = ["Content to be added"]
                else:
                    content[key] = "Content to be added"
        
        # Ensure body is always a list
        if not isinstance(content['body'], list):
            content['body'] = [str(content['body'])]
        
        if layout is None:
            layout = self.prs.slide_layouts[1]
        
        slide = self.prs.slides.add_slide(layout)
        slide_dict = {
            'content': content,
            'slide_object': slide
        }
        
        self.slides.append(slide_dict)
        self.total_slides += 1
        
        # Apply content immediately
        self.apply_content_to_slide(slide_dict)
        
        # Save after each slide
        self.save_version()
        
    def apply_content_to_slide(self, slide_dict):
        slide = slide_dict['slide_object']
        content = slide_dict['content']
        
        # Apply title
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            if 'title' in content:
                slide.shapes.title.text = content['title']
        
        # Apply body content
        if 'body' in content:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape != slide.shapes.title:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        if isinstance(content['body'], list):
                            # Handle bullet points
                            for i, bullet_point in enumerate(content['body']):
                                if i == 0:
                                    p.text = bullet_point
                                else:
                                    p = text_frame.add_paragraph()
                                    p.text = bullet_point
                                p.level = 0
                        else:
                            p.text = content['body']
        
        # Apply speaker notes
        if 'notes' in content:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content['notes']
    
    def save_version(self):
        """Save a versioned copy of the presentation"""
        self.version += 1
        timestamp = time.strftime("%Y%m%d-%H%M%S")
        version_filename = f"{self.filename.rsplit('.', 1)[0]}_v{self.version}_{timestamp}.pptx"
        
        try:
            # Save the current version
            self.prs.save(version_filename)
            # Also save/update the main file
            self.prs.save(self.filename)
            print(f"\nSaved presentation version {self.version} to {version_filename}")
            print(f"Updated main file: {self.filename}")
        except Exception as e:
            print(f"Error saving presentation: {str(e)}")
    
    def is_complete(self):
        return self.total_slides >= self.target_slides
    
    def get_progress(self):
        return f"Current progress: {self.total_slides}/{self.target_slides} slides (Version {self.version})"

presentation_progress = PresentationProgress()

def create_messages(system_message, user_message):
    return [
        {"role": "system", "content": system_message},
        {"role": "user", "content": f"{user_message}\n\nCurrent progress: {presentation_progress.get_progress()}"}
    ]

class CustomAssistantAgent(autogen.AssistantAgent):
    def generate_reply(self, sender=None, messages=None):
        if messages is None:
            messages = self._oai_messages
            
        if not messages:
            return None
            
        system_message = self.system_message
        
        last_message = messages[-1]
        if isinstance(last_message, dict):
            last_message_content = last_message.get('content', '')
        else:
            last_message_content = str(last_message)
        
        api_messages = create_messages(system_message, last_message_content)
        response = llm.create_completion(api_messages)
        
        try:
            reply = response['choices'][0]['message']['content']
            if self.name == "content_writer":
                try:
                    # Try to find JSON content within the response
                    json_start = reply.find('{')
                    json_end = reply.rfind('}') + 1
                    if json_start != -1 and json_end != -1:
                        json_content = reply[json_start:json_end]
                        content_dict = json.loads(json_content)
                        
                        # Validate required keys
                        required_keys = ['title', 'body', 'notes']
                        if all(key in content_dict for key in required_keys):
                            presentation_progress.add_slide(content_dict)
                            print(f"\nAdded new slide: {content_dict['title']}")
                        else:
                            missing_keys = [key for key in required_keys if key not in content_dict]
                            print(f"Warning: Missing required keys in JSON: {missing_keys}")
                            # Create a default structure for missing keys
                            for key in missing_keys:
                                if key == 'body':
                                    content_dict[key] = ["Content to be added"]
                                else:
                                    content_dict[key] = "Content to be added"
                            presentation_progress.add_slide(content_dict)
                    else:
                        print("Warning: Could not find JSON content in response")
                        # Create a default slide
                        default_content = {
                            "title": "Content Processing Error",
                            "body": ["Content could not be properly formatted", 
                                   "Please review and update this slide"],
                            "notes": "This slide needs to be reviewed and updated with proper content."
                        }
                        presentation_progress.add_slide(default_content)
                except json.JSONDecodeError as e:
                    print(f"Warning: JSON parsing error: {str(e)}")
                    print("Response content:", reply)
                    # Create a default slide
                    default_content = {
                        "title": "Content Processing Error",
                        "body": ["Content could not be properly formatted", 
                               "Please review and update this slide"],
                        "notes": "This slide needs to be reviewed and updated with proper content."
                    }
                    presentation_progress.add_slide(default_content)
            return reply
        except KeyError:
            return "I apologize, but I encountered an error processing your request."

# Initialize agents
content_strategist = CustomAssistantAgent(
    name="content_strategist",
    system_message=content_strategist_config["system_message"]
)

slide_designer = CustomAssistantAgent(
    name="slide_designer",
    system_message=slide_designer_config["system_message"]
)

content_writer = CustomAssistantAgent(
    name="content_writer",
    system_message=content_writer_config["system_message"]
)

visualization_expert = CustomAssistantAgent(
    name="visualization_expert",
    system_message=visualization_expert_config["system_message"]
)

user_proxy = autogen.UserProxyAgent(
    name="user_proxy",
    human_input_mode="TERMINATE",
    max_consecutive_auto_reply=10
)

groupchat = autogen.GroupChat(
    agents=[user_proxy, content_strategist, slide_designer, content_writer, visualization_expert],
    messages=[],
    max_round=50
)

manager = autogen.GroupChatManager(groupchat=groupchat)

initial_message = """Create a professional presentation about artificial intelligence trends in 2024.
Requirements:
1. At least 15 slides
2. Include executive summary
3. Cover key trends, market analysis, and future predictions
4. Include data visualizations and charts
5. Provide speaker notes for each slide

Please start by creating a detailed outline and design guidelines."""

# Generate presentation with continuous saves
while not presentation_progress.is_complete():
    if presentation_progress.total_slides == 0:
        user_proxy.initiate_chat(manager, message=initial_message)
    else:
        next_slide_prompt = f"""Create the next slide content and design. 
Current progress: {presentation_progress.get_progress()}
Ensure content is clear, impactful, and visually appealing."""
        user_proxy.initiate_chat(manager, message=next_slide_prompt)

print(f"\nFinal presentation statistics:\n{presentation_progress.get_progress()}")
print(f"Final presentation saved as: {presentation_progress.filename}")

